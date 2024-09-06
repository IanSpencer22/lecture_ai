import os
import tkinter as tk
from tkinter import filedialog, scrolledtext
import speech_recognition as sr
from moviepy.editor import VideoFileClip
import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from collections import Counter
from pathlib import Path
import openai
import PyPDF2
from pptx import Presentation
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Initialize OpenAI client
openai.api_key = os.getenv("OPENAI_API_KEY")

# Global variable to store the uploaded file content
file_content = ""

def transcribe_video(video_path):
    video = VideoFileClip(video_path)
    audio_path = "temp_audio.wav"
    video.audio.write_audiofile(audio_path)

    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        audio = recognizer.record(source)
        try:
            text = recognizer.recognize_google(audio)
        except sr.UnknownValueError:
            text = "Audio not clear or no recognizable speech found."
        except sr.RequestError as e:
            text = f"Could not request results from Google Speech Recognition service; {e}"

    return text

def read_text_file(file_path):
    with open(file_path, 'r') as file:
        text = file.read()
    return text

def read_pdf_file(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfFileReader(file)
        for page_num in range(reader.numPages):
            page = reader.getPage(page_num)
            text += page.extract_text()
    return text

def read_ppt_file(file_path):
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def analyze_content(text):
    words = word_tokenize(text)
    stop_words = set(stopwords.words('english'))
    filtered_words = [word for word in words if word.isalnum() and word.lower() not in stop_words]
    word_freq = Counter(filtered_words)
    common_words = word_freq.most_common(10)
    return common_words

def generate_ai_response(prompt, context):
    response = openai.chat.completions.create(
        messages=[
            {"role": "system", "content": "You are a helpful assistant. Do not use markdown formatting."},
            {"role": "user", "content": f"Context: {context}\n\nQuestion: {prompt}"}
        ],
        model="gpt-4o-mini",
    )
    return response.choices[0].message.content.strip()

def send_message():
    user_message = user_input.get()
    if user_message.strip():
        chat_history.config(state=tk.NORMAL)
        chat_history.insert(tk.END, f"You: {user_message}\n", "user")
        chat_history.config(state=tk.DISABLED)
        user_input.set("")

        ai_response = generate_ai_response(user_message, file_content)
        chat_history.config(state=tk.NORMAL)
        chat_history.insert(tk.END, f"AI: {ai_response}\n", "ai")
        chat_history.config(state=tk.DISABLED)
        chat_history.yview(tk.END)

def upload_file():
    global file_content
    file_path = filedialog.askopenfilename(filetypes=[("All supported files", "*.mp4;*.mkv;*.avi;*.txt;*.pdf;*.pptx")])
    if file_path:
        if file_path.endswith(('.mp4', '.mkv', '.avi')):
            file_content = transcribe_video(file_path)
        elif file_path.endswith('.txt'):
            file_content = read_text_file(file_path)
        elif file_path.endswith('.pdf'):
            file_content = read_pdf_file(file_path)
        elif file_path.endswith('.pptx'):
            file_content = read_ppt_file(file_path)
        else:
            file_content = "Unsupported file format."

        chat_history.config(state=tk.NORMAL)
        chat_history.insert(tk.END, f"File content loaded. You can now ask questions based on the uploaded content.\n")
        chat_history.config(state=tk.DISABLED)
        chat_history.yview(tk.END)

# GUI setup
root = tk.Tk()
root.title("Lecture AI Generator")

# Chat history display
chat_history = scrolledtext.ScrolledText(root, wrap=tk.WORD, state=tk.DISABLED)
chat_history.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)

# Add tags for styling
chat_history.tag_config("user", foreground="blue", font=("Helvetica", 12, "bold"))
chat_history.tag_config("ai", foreground="black", font=("Helvetica", 12))

# User input field
user_input = tk.StringVar()
input_frame = tk.Frame(root)
input_frame.pack(padx=10, pady=10, fill=tk.X)

input_entry = tk.Entry(input_frame, textvariable=user_input)
input_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))

send_button = tk.Button(input_frame, text="Send", command=send_message)
send_button.pack(side=tk.RIGHT)

# Upload file button
upload_button = tk.Button(root, text="Upload File", command=upload_file)
upload_button.pack(pady=10)

root.mainloop()